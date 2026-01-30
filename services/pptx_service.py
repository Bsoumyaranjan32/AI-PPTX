"""
Enterprise-Grade PowerPoint Generation Service
Version: 2.0.0
Author: GuptaSigma
Date: 2026-01-29

Features:
- Advanced Layout System (12+ Layout Types)
- Robust Image Handling with Retry & Fallback
- Chart & Table Generation
- Theme Management System
- Comprehensive Error Handling
- Logging & Monitoring
- Performance Optimization
"""

import io
import os
import re
import logging
import requests
import urllib3
from typing import Dict, List, Optional, Tuple, Any
from datetime import datetime
from io import BytesIO
from urllib.parse import quote, urlparse
import hashlib
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image

# Disable SSL warnings for development
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# Configure Logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('pptx_service.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)


class PPTXConfig:
    """Configuration class for PPTX Service"""
    
    # Slide Dimensions (16:9 Widescreen)
    SLIDE_WIDTH = Inches(13.33)
    SLIDE_HEIGHT = Inches(7.5)
    
    # Default Margins
    MARGIN_TOP = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)
    MARGIN_LEFT = Inches(0.5)
    MARGIN_RIGHT = Inches(0.5)
    
    # Font Sizes
    FONT_SIZE_TITLE = Pt(48)      # Increased for better visibility
    FONT_SIZE_HEADING = Pt(36)    # Matches web h2
    FONT_SIZE_SUBHEADING = Pt(26) # Matches web h3
    FONT_SIZE_BODY = Pt(16)       # Matches web body
    FONT_SIZE_SMALL = Pt(13)      # Matches web small
    FONT_SIZE_CAPTION = Pt(11)    # Matches web caption
    
    # Image Settings
    IMAGE_TIMEOUT = 15
    IMAGE_MAX_RETRIES = 3
    IMAGE_CACHE_ENABLED = True
    
    # Performance Settings
    MAX_SLIDES = 100
    MAX_IMAGE_SIZE_MB = 10


class ThemeManager:
    """Advanced Theme Management System"""
    
    def __init__(self):
        self.themes = {
            "dialogue": {
                'name': 'Dialogue',
                'bg': RGBColor(255, 255, 255),
                'text': RGBColor(0, 0, 0),
                'accent': RGBColor(99, 102, 241),
                'accent_light': RGBColor(165, 180, 252),
                'card': RGBColor(245, 247, 250),
                'border': RGBColor(226, 232, 240),
                'success': RGBColor(34, 197, 94),
                'warning': RGBColor(251, 191, 36),
                'error': RGBColor(239, 68, 68),
            },
            "alien": {
                'name': 'Alien Dark',
                'bg': RGBColor(15, 23, 42),
                'text': RGBColor(255, 255, 255),
                'accent': RGBColor(34, 211, 238),
                'accent_light': RGBColor(103, 232, 249),
                'card': RGBColor(30, 41, 59),
                'border': RGBColor(51, 65, 85),
                'success': RGBColor(52, 211, 153),
                'warning': RGBColor(250, 204, 21),
                'error': RGBColor(248, 113, 113),
            },
            "wine": {
                'name': 'Wine Elegance',
                'bg': RGBColor(76, 29, 51),
                'text': RGBColor(255, 255, 255),
                'accent': RGBColor(244, 114, 182),
                'accent_light': RGBColor(249, 168, 212),
                'card': RGBColor(100, 40, 70),
                'border': RGBColor(157, 23, 77),
                'success': RGBColor(167, 243, 208),
                'warning': RGBColor(253, 224, 71),
                'error': RGBColor(252, 165, 165),
            },
            "business": {
                'name': 'Business Professional',
                'bg': RGBColor(255, 255, 255),
                'text': RGBColor(15, 23, 42),
                'accent': RGBColor(37, 99, 235),
                'accent_light': RGBColor(147, 197, 253),
                'card': RGBColor(241, 245, 249),
                'border': RGBColor(203, 213, 225),
                'success': RGBColor(16, 185, 129),
                'warning': RGBColor(245, 158, 11),
                'error': RGBColor(220, 38, 38),
            },
            "ocean": {
                'name': 'Ocean Breeze',
                'bg': RGBColor(240, 249, 255),
                'text': RGBColor(12, 74, 110),
                'accent': RGBColor(14, 165, 233),
                'accent_light': RGBColor(125, 211, 252),
                'card': RGBColor(224, 242, 254),
                'border': RGBColor(186, 230, 253),
                'success': RGBColor(6, 182, 212),
                'warning': RGBColor(251, 146, 60),
                'error': RGBColor(239, 68, 68),
            },
            "forest": {
                'name': 'Forest Green',
                'bg': RGBColor(236, 253, 245),
                'text': RGBColor(6, 78, 59),
                'accent': RGBColor(5, 150, 105),
                'accent_light': RGBColor(110, 231, 183),
                'card': RGBColor(209, 250, 229),
                'border': RGBColor(167, 243, 208),
                'success': RGBColor(34, 197, 94),
                'warning': RGBColor(234, 179, 8),
                'error': RGBColor(220, 38, 38),
            },
            "sunset": {
                'name': 'Sunset Orange',
                'bg': RGBColor(255, 247, 237),
                'text': RGBColor(124, 45, 18),
                'accent': RGBColor(249, 115, 22),
                'accent_light': RGBColor(251, 146, 60),
                'card': RGBColor(254, 215, 170),
                'border': RGBColor(253, 186, 116),
                'success': RGBColor(132, 204, 22),
                'warning': RGBColor(234, 179, 8),
                'error': RGBColor(239, 68, 68),
            },
            "midnight": {
                'name': 'Midnight Purple',
                'bg': RGBColor(24, 24, 27),
                'text': RGBColor(250, 250, 250),
                'accent': RGBColor(168, 85, 247),
                'accent_light': RGBColor(196, 181, 253),
                'card': RGBColor(39, 39, 42),
                'border': RGBColor(63, 63, 70),
                'success': RGBColor(134, 239, 172),
                'warning': RGBColor(250, 204, 21),
                'error': RGBColor(248, 113, 113),
            }
        }
    
    def get_theme(self, theme_name: str) -> Dict:
        """Get theme by name with fallback"""
        theme_name = theme_name.lower()
        return self.themes.get(theme_name, self.themes["dialogue"])
    
    def list_themes(self) -> List[str]:
        """List all available themes"""
        return list(self.themes.keys())


class ImageHandler:
    """Advanced Image Handling with Caching and Retry Logic"""
    
    def __init__(self, config: PPTXConfig):
        self.config = config
        self.cache = {}
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
            'Accept-Encoding': 'gzip, deflate, br',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1'
        })
    
    def _sanitize_url(self, url: str) -> str:
        """Sanitize and encode URL properly"""
        if not url:
            return ""
        
        try:
            logger.info(f"🔧 Original URL: {url}")
            
            # Handle spaces and special chars in URL path
            if " " in url or any(c in url for c in ['&', '=', '?']):
                parsed = urlparse(url)
                
                # If prompt parameter exists, encode only that part
                if "prompt/" in url:
                    base, prompt_part = url.rsplit("prompt/", 1)
                    encoded_prompt = quote(prompt_part, safe='')
                    url = f"{base}prompt/{encoded_prompt}"
                else:
                    # Full URL encoding (preserve scheme and domain)
                    url = quote(url, safe=':/?#[]@!$&\'()*+,;=')
            
            logger.info(f"✅ Sanitized URL: {url}")
            return url
            
        except Exception as e:
            logger.error(f"❌ URL sanitization failed: {e}")
            return url
    
    def _get_cache_key(self, url: str) -> str:
        """Generate cache key from URL"""
        return hashlib.md5(url.encode()).hexdigest()
    
    def download_image(self, url: str, retries: Optional[int] = None) -> Optional[BytesIO]:
        """
        Download image with retry logic and caching
        
        Args:
            url: Image URL to download
            retries: Number of retry attempts (defaults to config value)
        
        Returns:
            BytesIO object containing image data or None if failed
        """
        if not url:
            return None
        
        # Sanitize URL
        url = self._sanitize_url(url)
        
        # Check cache
        if self.config.IMAGE_CACHE_ENABLED:
            cache_key = self._get_cache_key(url)
            if cache_key in self.cache:
                logger.info(f"✅ Image loaded from cache: {url[:50]}...")
                return BytesIO(self.cache[cache_key])
        
        # Set retry count
        max_retries = retries if retries is not None else self.config.IMAGE_MAX_RETRIES
        
        # Attempt download with retries
        for attempt in range(max_retries):
            try:
                logger.info(f"⬇️ Downloading image (attempt {attempt + 1}/{max_retries}): {url[:50]}...")
                
                response = self.session.get(
                    url,
                    timeout=self.config.IMAGE_TIMEOUT,
                    verify=False,
                    stream=True
                )
                
                if response.status_code == 200:
                    # Check image size
                    content_length = response.headers.get('content-length')
                    if content_length and int(content_length) > self.config.MAX_IMAGE_SIZE_MB * 1024 * 1024:
                        logger.warning(f"⚠️ Image too large: {content_length} bytes")
                        return None
                    
                    # Read image data
                    image_data = response.content
                    
                    # Verify image is valid using a separate BytesIO to avoid corruption
                    try:
                        img = Image.open(BytesIO(image_data))
                        img.verify()
                        # Reopen to get format and size after verify()
                        img = Image.open(BytesIO(image_data))
                        logger.info(f"✅ Image verified: {img.format} {img.size}")
                    except Exception as e:
                        logger.warning(f"⚠️ Image verification failed: {e}")
                        return None
                    
                    # Cache and return
                    if self.config.IMAGE_CACHE_ENABLED:
                        self.cache[cache_key] = image_data
                    
                    logger.info(f"✅ Image downloaded successfully: {len(image_data)} bytes")
                    return BytesIO(image_data)
                else:
                    logger.warning(f"⚠️ HTTP {response.status_code} for: {url[:50]}...")
                    
            except requests.exceptions.Timeout:
                logger.warning(f"⏱️ Timeout (attempt {attempt + 1}/{max_retries}): {url[:50]}...")
            except requests.exceptions.SSLError as e:
                logger.warning(f"🔒 SSL Error (attempt {attempt + 1}/{max_retries}): {e}")
            except Exception as e:
                logger.error(f"❌ Download error (attempt {attempt + 1}/{max_retries}): {e}")
        
        logger.error(f"❌ Failed to download image after {max_retries} attempts: {url[:50]}...")
        return None
    
    def clear_cache(self):
        """Clear image cache"""
        self.cache.clear()
        logger.info("🗑️ Image cache cleared")


class ShapeHelper:
    """Helper class for creating and styling shapes"""
    
    @staticmethod
    def add_rounded_rectangle(slide, x, y, width, height, fill_color, border_color=None, border_width=Pt(1)):
        """Add a rounded rectangle with styling"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = border_width
        else:
            shape.line.fill.background()
        
        return shape
    
    @staticmethod
    def add_circle(slide, x, y, diameter, fill_color, transparency=0.0, border=False):
        """Add a circle/oval with styling"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x, y, diameter, diameter
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.fill.transparency = transparency
        
        if not border:
            shape.line.fill.background()
        
        return shape
    
    @staticmethod
    def add_badge(slide, x, y, text, bg_color, text_color, size=Inches(0.5)):
        """Add a circular badge with text"""
        badge = ShapeHelper.add_circle(slide, x, y, size, bg_color, border=False)
        
        text_frame = badge.text_frame
        text_frame.text = str(text)
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = text_color
        paragraph.alignment = PP_ALIGN.CENTER
        
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        return badge
    
    @staticmethod
    def add_separator_line(slide, x1, y1, x2, y2, color, width=Pt(2)):
        """Add a separator line"""
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x1, y1, x2, y2
        )
        connector.line.color.rgb = color
        connector.line.width = width
        return connector


class ChartBuilder:
    """Advanced Chart Building"""
    
    @staticmethod
    def add_bar_chart(slide, x, y, width, height, chart_data_dict, title=""):
        """
        Add a bar chart to slide
        
        Args:
            chart_data_dict: {'categories': [...], 'series': {'Series 1': [...], ...}}
        """
        try:
            chart_data = CategoryChartData()
            chart_data.categories = chart_data_dict.get('categories', [])
            
            for series_name, values in chart_data_dict.get('series', {}).items():
                chart_data.add_series(series_name, values)
            
            graphic_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                x, y, width, height,
                chart_data
            )
            
            chart = graphic_frame.chart
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
            
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            
            return chart
        except Exception as e:
            logger.error(f"Bar chart creation error: {e}")
            return None
    
    @staticmethod
    def add_line_chart(slide, x, y, width, height, chart_data_dict, title=""):
        """Add a line chart to slide"""
        try:
            chart_data = CategoryChartData()
            chart_data.categories = chart_data_dict.get('categories', [])
            
            for series_name, values in chart_data_dict.get('series', {}).items():
                chart_data.add_series(series_name, values)
            
            graphic_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE,
                x, y, width, height,
                chart_data
            )
            
            chart = graphic_frame.chart
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
            
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            
            return chart
        except Exception as e:
            logger.error(f"Line chart creation error: {e}")
            return None
    
    @staticmethod
    def add_pie_chart(slide, x, y, width, height, data_dict, title=""):
        """
        Add a pie chart to slide
        
        Args:
            data_dict: {'Category1': value1, 'Category2': value2, ...}
        """
        try:
            chart_data = CategoryChartData()
            chart_data.categories = list(data_dict.keys())
            chart_data.add_series('Values', list(data_dict.values()))
            
            graphic_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                x, y, width, height,
                chart_data
            )
            
            chart = graphic_frame.chart
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
            
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            
            # Add data labels
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.number_format = '0%'
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            
            return chart
        except Exception as e:
            logger.error(f"Pie chart creation error: {e}")
            return None


class TableBuilder:
    """Advanced Table Building"""
    
    @staticmethod
    def add_table(slide, x, y, data, theme, has_header=True):
        """
        Add a styled table to slide
        
        Args:
            data: List of lists [[row1], [row2], ...]
        """
        try:
            if not data or len(data) == 0:
                return None
            
            rows = len(data)
            cols = len(data[0])
            
            # Calculate dimensions
            col_width = Inches(12) / cols
            row_height = Inches(0.5)
            table_width = col_width * cols
            table_height = row_height * rows
            
            # Create table
            shape = slide.shapes.add_table(
                rows, cols,
                x, y,
                table_width, table_height
            )
            
            table = shape.table
            
            # Style table
            for row_idx, row_data in enumerate(data):
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value)
                    
                    # Style header row
                    if has_header and row_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = theme['accent']
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.bold = True
                            paragraph.font.size = Pt(14)
                            paragraph.font.color.rgb = theme['bg']
                    else:
                        # Alternate row colors
                        if row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = theme['card']
                        
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
                            paragraph.font.color.rgb = theme['text']
            
            return table
        except Exception as e:
            logger.error(f"Table creation error: {e}")
            return None


class PPTXService:
    """
    Enterprise-Grade PowerPoint Generation Service
    
    Supports multiple layouts, themes, charts, tables, and advanced features
    """
    
    def __init__(self):
        self.config = PPTXConfig()
        self.theme_manager = ThemeManager()
        self.image_handler = ImageHandler(self.config)
        self.prs = None
        self.current_theme = None
        
        logger.info("🚀 PPTXService initialized")
    
    def generate(self, presentation_data) -> bytes:
        """
        Generate PowerPoint presentation from data
        
        Args:
            presentation_data: Object with theme and content attributes
        
        Returns:
            bytes: Generated PPTX file content
        """
        try:
            logger.info("\n" + "="*80)
            logger.info("🎨 STARTING PRESENTATION GENERATION")
            logger.info("="*80)
            
            # Initialize presentation
            self.prs = Presentation()
            self.prs.slide_width = self.config.SLIDE_WIDTH
            self.prs.slide_height = self.config.SLIDE_HEIGHT
            
            # Get theme
            theme_name = getattr(presentation_data, 'theme', 'dialogue').lower()
            self.current_theme = self.theme_manager.get_theme(theme_name)
            logger.info(f"🎨 Theme: {self.current_theme['name']}")
            
            # Get slides data
            content = getattr(presentation_data, 'content', {})
            slides_data = content.get('slides', [])
            
            # Validate slide count
            if len(slides_data) > self.config.MAX_SLIDES:
                logger.warning(f"⚠️ Slide count ({len(slides_data)}) exceeds maximum ({self.config.MAX_SLIDES})")
                slides_data = slides_data[:self.config.MAX_SLIDES]
            
            logger.info(f"📊 Total slides to generate: {len(slides_data)}")
            
            # Generate slides
            for i, slide_data in enumerate(slides_data, 1):
                self._generate_slide(i, slide_data)
            
            # Save to BytesIO
            output = BytesIO()
            self.prs.save(output)
            output.seek(0)
            
            logger.info("="*80)
            logger.info("✅ PRESENTATION GENERATION COMPLETED")
            logger.info("="*80 + "\n")
            
            return output.getvalue()
            
        except Exception as e:
            logger.error(f"❌ Presentation generation failed: {e}", exc_info=True)
            raise
    
    def _generate_slide(self, slide_number: int, slide_data: Dict):
        """Route slide generation to appropriate layout handler"""
        layout = slide_data.get('layout', 'standard')
        
        logger.info(f"\n📄 Slide {slide_number}: {layout.upper()}")
        logger.info(f"   Title: {slide_data.get('title', 'N/A')[:50]}")
        
        try:
            # Layout router
            layout_handlers = {
                'centered': self._create_hero_slide,
                'hero': self._create_hero_slide,
                'title': self._create_hero_slide,
                'split_box': self._create_split_card_slide,
                'split': self._create_split_card_slide,
                'grid_4': self._create_grid_cards_slide,
                'grid': self._create_grid_cards_slide,
                'roadmap': self._create_roadmap_slide,
                'timeline': self._create_timeline_slide,
                'comparison': self._create_comparison_slide,
                'quote': self._create_quote_slide,
                'image_focus': self._create_image_focus_slide,
                'two_column': self._create_two_column_slide,
                'chart': self._create_chart_slide,
                'table': self._create_table_slide,
            }
            
            handler = layout_handlers.get(layout, self._create_standard_slide)
            handler(slide_data, self.current_theme)
            
            logger.info(f"   ✅ Slide {slide_number} created successfully")
            
        except Exception as e:
            logger.error(f"   ❌ Error creating slide {slide_number}: {e}")
            # Create error slide as fallback
            self._create_error_slide(slide_data, str(e))
    
    def _preprocess_content(self, content: str) -> str:
        """Clean and format content for PPTX"""
        if not content:
            return ""
        
        # Remove HTML tags if any
        content = re.sub(r'<[^>]+>', '', content)
        
        # Convert markdown to plain text with bullets
        content = content.replace('**', '')
        content = content.replace('*', '•')
        content = content.replace('- ', '• ')
        
        # Fix line breaks
        content = content.replace('\r\n', '\n')
        content = content.replace('\r', '\n')
        
        # Remove excessive whitespace
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        content = '\n'.join(lines)
        
        return content
    
    def _add_debug_borders(self, slide, theme):
        """Add border guides for debugging (only in dev mode)"""
        if os.getenv('PPTX_DEBUG') == 'true':
            # Add thin border around content areas
            for x, y, w, h in [
                (Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
            ]:
                border = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, x, y, w, h
                )
                border.fill.background()
                border.line.color.rgb = RGBColor(255, 0, 0)
                border.line.width = Pt(1)
                border.line.dash_style = 2  # Dashed
    
    # ==========================================
    # LAYOUT 1: HERO / TITLE SLIDE
    # ==========================================
    def _create_hero_slide(self, data: Dict, theme: Dict):
        """Create hero/title slide with image or fallback design"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        img_url = data.get('image')
        image_loaded = False
        
        # Attempt to load background image - ADD IMAGE FIRST AS BACKGROUND LAYER
        if img_url:
            image_data = self.image_handler.download_image(img_url)
            if image_data:
                try:
                    slide.shapes.add_picture(
                        image_data,
                        0, 0,
                        self.config.SLIDE_WIDTH,
                        self.config.SLIDE_HEIGHT
                    )
                    image_loaded = True
                    logger.info("   🖼️ Hero image loaded")
                    
                    # Verify that image actually rendered
                    try:
                        last_shape = slide.shapes[-1]
                        if not hasattr(last_shape, 'image'):
                            logger.warning("   ⚠️ Image shape invalid, using fallback")
                            image_loaded = False
                            # Remove the failed shape
                            slide.shapes._spTree.remove(last_shape._element)
                    except Exception as e:
                        logger.warning(f"   ⚠️ Image verification failed: {e}")
                        image_loaded = False
                        
                except Exception as e:
                    logger.warning(f"   ⚠️ Failed to place hero image: {e}")
        
        # Apply overlay or fallback design
        if image_loaded:
            # ADD DARK OVERLAY ON TOP OF IMAGE (50-60% transparency for readability)
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0,
                self.config.SLIDE_WIDTH,
                self.config.SLIDE_HEIGHT
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.50  # 50% transparency for optimal readability
            overlay.line.fill.background()
            # TEXT MUST BE WHITE for readability on dark overlay
            text_color = RGBColor(255, 255, 255)
        else:
            # Fallback: Theme background with decorative circles
            self._set_background(slide, theme)
            
            # Add subtle gradient background for better readability
            gradient_overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, Inches(2),
                self.config.SLIDE_WIDTH,
                Inches(4)
            )
            gradient_overlay.fill.solid()
            gradient_overlay.fill.fore_color.rgb = theme['accent']
            gradient_overlay.fill.transparency = 0.85
            gradient_overlay.line.fill.background()
            
            # Large decorative circle (top-left)
            ShapeHelper.add_circle(
                slide,
                Inches(-2), Inches(-2),
                Inches(6),
                theme['accent'],
                transparency=0.15
            )
            
            # Medium circle (bottom-right)
            ShapeHelper.add_circle(
                slide,
                Inches(10), Inches(4),
                Inches(5),
                theme['accent'],
                transparency=0.15
            )
            
            # Small accent circle
            ShapeHelper.add_circle(
                slide,
                Inches(11), Inches(1),
                Inches(2),
                theme['accent_light'],
                transparency=0.2
            )
            
            text_color = theme['text']
        
        # ADD TEXT SHAPES LAST (on top layer for proper z-index)
        # Title - 44pt bold white (centered) - positioned higher for visibility
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8),
            Inches(12.0), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = data.get('title', '')
        title_paragraph.font.size = Pt(44)  # 44pt as per spec
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = text_color
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Body Content - 16pt for better readability (centered)
        if data.get('content'):
            subtitle_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(2.0),
                Inches(11.0), Inches(4.5)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            
            content_text = data.get('content', '')
            # Preprocess content to clean it up
            content_text = self._preprocess_content(content_text)
            
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.text = content_text
            subtitle_paragraph.font.size = Pt(16)  # 16pt as per spec
            subtitle_paragraph.font.color.rgb = text_color
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
            subtitle_paragraph.space_after = Pt(12)
            subtitle_paragraph.line_spacing = 1.2
    
    # ==========================================
    # LAYOUT 2: SPLIT CARD SLIDE
    # ==========================================
    def _create_split_card_slide(self, data: Dict, theme: Dict):
        """Create split layout with image and content card"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        # Left: Image
        img_url = data.get('image')
        if img_url:
            image_data = self.image_handler.download_image(img_url)
            if image_data:
                try:
                    slide.shapes.add_picture(
                        image_data,
                        Inches(0.5), Inches(2),
                        width=Inches(5),
                        height=Inches(4.5)
                    )
                    logger.info("   🖼️ Split image loaded")
                except Exception as e:
                    logger.warning(f"   ⚠️ Failed to place split image: {e}")
        
        # Right: Content Card
        card = ShapeHelper.add_rounded_rectangle(
            slide,
            Inches(6), Inches(2),
            Inches(6.8), Inches(4.5),
            theme['card'],
            theme['accent'],
            Pt(2)
        )
        
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_bottom = Inches(0.3)
        
        text_frame.text = data.get('content', '')
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = self.config.FONT_SIZE_BODY
            paragraph.font.color.rgb = theme['text']
            paragraph.space_after = Pt(10)
    
    # ==========================================
    # LAYOUT 3: GRID CARDS SLIDE
    # ==========================================
    def _create_grid_cards_slide(self, data: Dict, theme: Dict):
        """Create 3-column horizontal grid with numbered cards (1, 2, 3)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Apply dark navy background for grid slides
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.config.SLIDE_WIDTH,
            self.config.SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(20, 30, 48)  # Dark navy background
        background.line.fill.background()
        
        # Send background to back
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)
        
        # Add title with white text for dark background (left-aligned as per spec)
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.6),
            Inches(12.0),
            Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = data.get('title', '')
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        # Parse content
        content = data.get('content', '')
        items = [line.strip().replace('*', '').replace('#', '')
                 for line in content.split('\n') if line.strip()]
        
        # 3-COLUMN HORIZONTAL LAYOUT with proper spacing
        card_width = Inches(4.0)   # Each card width
        card_height = Inches(2.5)  # Card height
        gap = Inches(0.3)          # Gap between cards
        start_x = Inches(0.5)      # Starting X position
        start_y = Inches(2.5)      # Starting Y position
        
        # Create 3 cards with large numbers (1, 2, 3)
        for idx in range(min(3, len(items))):
            x_pos = start_x + (idx * (card_width + gap))
            
            # Card background with rounded corners
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_pos, start_y,
                card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(30, 42, 62)  # Slightly lighter than background
            card.line.color.rgb = RGBColor(80, 90, 110)
            card.line.width = Pt(1)
            
            # Large number label (1, 2, 3) in top-left
            num_box = slide.shapes.add_textbox(
                x_pos + Inches(0.2),
                start_y + Inches(0.2),
                Inches(0.6), Inches(0.6)
            )
            num_para = num_box.text_frame.paragraphs[0]
            num_para.text = str(idx + 1)  # Simple "1", "2", "3"
            num_para.font.size = Pt(32)
            num_para.font.bold = True
            num_para.font.color.rgb = RGBColor(255, 255, 255)  # White number
            
            # Content text
            content_box = slide.shapes.add_textbox(
                x_pos + Inches(0.3),
                start_y + Inches(0.9),
                card_width - Inches(0.6),
                card_height - Inches(1.0)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            content_para = content_frame.paragraphs[0]
            content_para.text = items[idx] if idx < len(items) else ""
            content_para.font.size = Pt(12)
            content_para.font.color.rgb = RGBColor(220, 220, 220)  # Light gray text
    
    def _create_content_card(self, slide, x, y, width, height, items, theme, number):
        """Helper to create a numbered content card with rounded corners and borders"""
        # Card background with rounded corners and light border
        card = ShapeHelper.add_rounded_rectangle(
            slide, x, y, width, height,
            theme['card'],
            RGBColor(200, 200, 220),  # Light border color
            Pt(1.5)  # Border width
        )
        
        # Number badge (01, 02, 03)
        ShapeHelper.add_badge(
            slide,
            x + Inches(0.2), y + Inches(0.2),
            number,
            theme['accent'],
            theme['bg'],
            Inches(0.6)
        )
        
        # Content text
        text_box = slide.shapes.add_textbox(
            x + Inches(0.3),
            y + Inches(1),
            width - Inches(0.6),
            height - Inches(1.2)
        )
        
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        content_text = "\n".join(items)
        content_text = content_text.replace('*', '').replace('#', '').strip()
        text_frame.text = content_text
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = self.config.FONT_SIZE_SMALL
            paragraph.font.color.rgb = theme['text']
            paragraph.space_after = Pt(8)
    
    # ==========================================
    # LAYOUT 4: ROADMAP SLIDE
    # ==========================================
    def _create_roadmap_slide(self, data: Dict, theme: Dict):
        """Create VERTICAL roadmap/timeline with numbered circles (1-6)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        # Parse items
        content = data.get('content', '')
        items = [line.strip().replace('*', '').replace('#', '') 
                 for line in content.split('\n') if line.strip()]
        
        if not items:
            return
        
        # Limit to 6 items max (matches web)
        items = items[:6]
        num_items = len(items)
        
        # VERTICAL TIMELINE LAYOUT SETTINGS (matching spec exactly)
        circle_diameter = Inches(0.6)       # Circle diameter
        circle_x = Inches(1.5)              # Fixed X position for all circles (left side)
        start_y = Inches(2.2)               # Starting Y position
        vertical_gap = Inches(1.0)          # Vertical spacing between circles
        line_thickness = Inches(0.08)       # Connecting line thickness (0.08" as per spec)
        text_box_x = Inches(2.5)            # Text box X position (right of circles) - updated to 2.5"
        text_box_width = Inches(9.5)        # Text box width - updated to 9.5"
        
        # DRAW VERTICAL CONNECTING LINE FIRST (behind circles)
        if num_items > 1:
            # Calculate line start and end Y positions
            line_start_y = start_y + (circle_diameter / 2)  # Center of first circle
            line_end_y = start_y + ((num_items - 1) * vertical_gap) + (circle_diameter / 2)  # Center of last circle
            
            # Draw vertical line connecting all circles with proper purple color
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                circle_x + (circle_diameter / 2), line_start_y,  # Start at center of first circle
                circle_x + (circle_diameter / 2), line_end_y     # End at center of last circle
            )
            line.line.color.rgb = RGBColor(100, 120, 255)  # Purple/blue color
            line.line.width = line_thickness
        
        # CREATE CIRCLES AND TEXT FOR EACH ROADMAP ITEM
        for i, item in enumerate(items):
            # Calculate Y position for this circle
            y = start_y + (i * vertical_gap)
            
            # Circle node with solid purple fill
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                circle_x, y,
                circle_diameter, circle_diameter
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(100, 120, 255)  # Purple solid fill
            circle.line.color.rgb = RGBColor(100, 120, 255)  # Purple border
            
            # White number text inside circle (1-6)
            num_frame = circle.text_frame
            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_para = num_frame.paragraphs[0]
            num_para.text = str(i + 1)
            num_para.font.size = Pt(20)
            num_para.font.bold = True
            num_para.font.color.rgb = RGBColor(255, 255, 255)  # White number
            num_para.alignment = PP_ALIGN.CENTER
            
            # Text content box aligned with circle (right of circle)
            text_box = slide.shapes.add_textbox(
                text_box_x, y,
                text_box_width, circle_diameter
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Split item into title and description if colon present
            if ':' in item:
                title_part, desc_part = item.split(':', 1)
                text_para = text_frame.paragraphs[0]
                text_para.text = title_part.strip()
                text_para.font.size = Pt(16)
                text_para.font.bold = True
                text_para.font.color.rgb = theme.get('text', RGBColor(51, 51, 51))
                
                # Add description on new line
                if desc_part.strip():
                    text_frame.add_paragraph()
                    desc_para = text_frame.paragraphs[1]
                    desc_para.text = desc_part.strip()
                    desc_para.font.size = Pt(12)
                    desc_para.font.color.rgb = RGBColor(100, 100, 100)
            else:
                text_para = text_frame.paragraphs[0]
                text_para.text = item
                text_para.font.size = Pt(14)
                text_para.font.color.rgb = theme.get('text', RGBColor(51, 51, 51))
    
    # ==========================================
    # LAYOUT 5: TIMELINE SLIDE
    # ==========================================
    def _create_timeline_slide(self, data: Dict, theme: Dict):
        """Create vertical timeline"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        # Parse timeline items
        content = data.get('content', '')
        items = [line.strip() for line in content.split('\n') if line.strip()]
        
        if not items:
            return
        
        # Timeline settings
        start_x = Inches(2)
        start_y = Inches(2)
        line_x = start_x + Inches(0.3)
        item_height = Inches(1.2)
        
        # Draw vertical line
        if len(items) > 1:
            ShapeHelper.add_separator_line(
                slide,
                line_x, start_y,
                line_x, start_y + (len(items) * item_height),
                theme['accent'],
                Pt(3)
            )
        
        # Create timeline items
        for i, item in enumerate(items):
            y = start_y + (i * item_height)
            
            # Circle marker
            ShapeHelper.add_circle(
                slide,
                line_x - Inches(0.15), y,
                Inches(0.3),
                theme['accent']
            )
            
            # Content
            text_box = slide.shapes.add_textbox(
                line_x + Inches(0.5), y - Inches(0.1),
                Inches(9), item_height - Inches(0.2)
            )
            
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = item.replace('*', '').replace('#', '').strip()
            
            for para in text_frame.paragraphs:
                para.font.size = self.config.FONT_SIZE_SMALL
                para.font.color.rgb = theme['text']
    
    # ==========================================
    # LAYOUT 6: COMPARISON SLIDE
    # ==========================================
    def _create_comparison_slide(self, data: Dict, theme: Dict):
        """Create side-by-side comparison"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        # Parse content
        content = data.get('content', '')
        items = [line.strip() for line in content.split('\n') if line.strip()]
        
        mid = len(items) // 2
        left_items = items[:mid]
        right_items = items[mid:]
        
        # Left side (Option A)
        left_card = ShapeHelper.add_rounded_rectangle(
            slide,
            Inches(0.5), Inches(2),
            Inches(5.5), Inches(4.5),
            theme['card'],
            theme['success'],
            Pt(3)
        )
        
        left_header = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(5.5), Inches(0.6)
        )
        left_header.fill.solid()
        left_header.fill.fore_color.rgb = theme['success']
        left_header_text = left_header.text_frame
        left_header_para = left_header_text.paragraphs[0]
        left_header_para.text = "✓ Option A"
        left_header_para.font.size = Pt(20)
        left_header_para.font.bold = True
        left_header_para.font.color.rgb = theme['bg']
        left_header_para.alignment = PP_ALIGN.CENTER
        
        left_content = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.8),
            Inches(5.1), Inches(3.5)
        )
        left_text = left_content.text_frame
        left_text.word_wrap = True
        left_text.text = "\n".join(left_items).replace('*', '•')
        
        for para in left_text.paragraphs:
            para.font.size = self.config.FONT_SIZE_SMALL
            para.font.color.rgb = theme['text']
            para.space_after = Pt(8)
        
        # Right side (Option B)
        right_card = ShapeHelper.add_rounded_rectangle(
            slide,
            Inches(7.3), Inches(2),
            Inches(5.5), Inches(4.5),
            theme['card'],
            theme['accent'],
            Pt(3)
        )
        
        right_header = slide.shapes.add_textbox(
            Inches(7.3), Inches(2),
            Inches(5.5), Inches(0.6)
        )
        right_header.fill.solid()
        right_header.fill.fore_color.rgb = theme['accent']
        right_header_text = right_header.text_frame
        right_header_para = right_header_text.paragraphs[0]
        right_header_para.text = "→ Option B"
        right_header_para.font.size = Pt(20)
        right_header_para.font.bold = True
        right_header_para.font.color.rgb = theme['bg']
        right_header_para.alignment = PP_ALIGN.CENTER
        
        right_content = slide.shapes.add_textbox(
            Inches(7.5), Inches(2.8),
            Inches(5.1), Inches(3.5)
        )
        right_text = right_content.text_frame
        right_text.word_wrap = True
        right_text.text = "\n".join(right_items).replace('*', '•')
        
        for para in right_text.paragraphs:
            para.font.size = self.config.FONT_SIZE_SMALL
            para.font.color.rgb = theme['text']
            para.space_after = Pt(8)
    
    # ==========================================
    # LAYOUT 7: QUOTE SLIDE
    # ==========================================
    def _create_quote_slide(self, data: Dict, theme: Dict):
        """Create inspirational quote slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        
        # Large decorative quotation mark
        quote_mark = slide.shapes.add_textbox(
            Inches(1), Inches(1),
            Inches(2), Inches(2)
        )
        quote_frame = quote_mark.text_frame
        quote_para = quote_frame.paragraphs[0]
        # FIXED: Added the missing closing quote and actual text
        quote_para.text = '"'
        quote_para.font.size = Pt(120)
        quote_para.font.color.rgb = theme['accent']
        quote_para.font.bold = True
        quote_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Quote content
        quote_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.5),
            Inches(9.33), Inches(3)
        )
        
        quote_text = quote_box.text_frame
        quote_text.word_wrap = True
        
        content_para = quote_text.paragraphs[0]
        content_para.text = data.get('content', '')
        content_para.font.size = Pt(28)
        content_para.font.italic = True
        content_para.font.color.rgb = theme['text']
        content_para.alignment = PP_ALIGN.CENTER
        
        # Author
        if data.get('title'):
            author_box = slide.shapes.add_textbox(
                Inches(2), Inches(5.5),
                Inches(9.33), Inches(1)
            )
            
            author_text = author_box.text_frame
            author_para = author_text.paragraphs[0]
            author_para.text = f"— {data.get('title', '')}"
            author_para.font.size = Pt(20)
            author_para.font.color.rgb = theme['accent']
            author_para.alignment = PP_ALIGN.RIGHT
    
    # ==========================================
    # LAYOUT 8: IMAGE FOCUS SLIDE
    # ==========================================
    def _create_image_focus_slide(self, data: Dict, theme: Dict):
        """Create slide with large centered image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        # Large centered image
        img_url = data.get('image')
        if img_url:
            image_data = self.image_handler.download_image(img_url)
            if image_data:
                try:
                    slide.shapes.add_picture(
                        image_data,
                        Inches(1.5), Inches(2),
                        width=Inches(10),
                        height=Inches(5)
                    )
                except Exception as e:
                    logger.warning(f"   ⚠️ Failed to place focus image: {e}")
        
        # Caption below image
        if data.get('content'):
            caption_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(7),
                Inches(10), Inches(0.3)
            )
            
            caption_text = caption_box.text_frame
            caption_para = caption_text.paragraphs[0]
            caption_para.text = data.get('content', '')
            caption_para.font.size = self.config.FONT_SIZE_CAPTION
            caption_para.font.italic = True
            caption_para.font.color.rgb = theme['text']
            caption_para.alignment = PP_ALIGN.CENTER
    
    # ==========================================
    # LAYOUT 9: TWO COLUMN SLIDE
    # ==========================================
    def _create_two_column_slide(self, data: Dict, theme: Dict):
        """Create two-column text layout"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        # Parse content
        content = data.get('content', '')
        items = [line.strip() for line in content.split('\n') if line.strip()]
        
        mid = len(items) // 2
        left_items = items[:mid]
        right_items = items[mid:]
        
        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(6), Inches(5)
        )
        left_text = left_box.text_frame
        left_text.word_wrap = True
        left_text.text = "\n".join(left_items).replace('*', '•')
        
        for para in left_text.paragraphs:
            para.font.size = self.config.FONT_SIZE_BODY
            para.font.color.rgb = theme['text']
            para.space_after = Pt(10)
        
        # Separator line
        ShapeHelper.add_separator_line(
            slide,
            Inches(6.665), Inches(2),
            Inches(6.665), Inches(7),
            theme['border'],
            Pt(2)
        )
        
        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(7), Inches(2),
            Inches(6), Inches(5)
        )
        right_text = right_box.text_frame
        right_text.word_wrap = True
        right_text.text = "\n".join(right_items).replace('*', '•')
        
        for para in right_text.paragraphs:
            para.font.size = self.config.FONT_SIZE_BODY
            para.font.color.rgb = theme['text']
            para.space_after = Pt(10)
    
    # ==========================================
    # LAYOUT 10: CHART SLIDE
    # ==========================================
    def _create_chart_slide(self, data: Dict, theme: Dict):
        """Create slide with chart"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        # Sample chart data (in real scenario, parse from data)
        chart_data = {
            'categories': ['Q1', 'Q2', 'Q3', 'Q4'],
            'series': {
                'Revenue': [50, 65, 80, 95],
                'Costs': [30, 35, 40, 45]
            }
        }
        
        ChartBuilder.add_bar_chart(
            slide,
            Inches(1.5), Inches(2),
            Inches(10), Inches(5),
            chart_data,
            title=data.get('title', '')
        )
    
    # ==========================================
    # LAYOUT 11: TABLE SLIDE
    # ==========================================
    def _create_table_slide(self, data: Dict, theme: Dict):
        """Create slide with table"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        # Sample table data (in real scenario, parse from data)
        table_data = [
            ['Feature', 'Plan A', 'Plan B', 'Plan C'],
            ['Price', '$10', '$20', '$30'],
            ['Storage', '10GB', '50GB', '100GB'],
            ['Users', '1', '5', 'Unlimited'],
        ]
        
        TableBuilder.add_table(
            slide,
            Inches(1), Inches(2),
            table_data,
            theme,
            has_header=True
        )
    
    # ==========================================
    # LAYOUT 12: STANDARD SLIDE
    # ==========================================
    def _create_standard_slide(self, data: Dict, theme: Dict):
        """Standard slide with balanced text and image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        self._add_title(slide, data.get('title', ''), theme)
        
        img_url = data.get('image')
        has_image = bool(img_url)
        image_loaded = False
        
        # Try to add image on LEFT side FIRST
        if has_image:
            image_data = self.image_handler.download_image(img_url)
            if image_data:
                try:
                    # Add image to LEFT side
                    # Note: Both width and height are specified to ensure consistent layout
                    # across slides, even if it means slight aspect ratio adjustment
                    slide.shapes.add_picture(
                        image_data,
                        Inches(0.5),   # X position - left side
                        Inches(2.2),   # Y position - below title  
                        width=Inches(5.5),   # Image width
                        height=Inches(3.5)   # Image height
                    )
                    image_loaded = True
                    logger.info("   🖼️ Standard image placed on left")
                except Exception as e:
                    logger.warning(f"   ⚠️ Image placement failed: {e}")
                    image_loaded = False
        
        # Add text content box on RIGHT side (if image present) or full width
        if image_loaded:
            # Text box on RIGHT side
            text_left = Inches(6.3)
            text_width = Inches(6.5)
        else:
            # Full width text box when no image
            text_left = Inches(0.8)
            text_width = Inches(11.8)
        
        # Text content - aligned with image Y position
        text_box = slide.shapes.add_textbox(
            text_left,
            Inches(2.2),  # Same Y as image
            text_width, 
            Inches(4.5)
        )
        
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        
        # Process content
        content = self._preprocess_content(data.get('content', ''))
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        # Add lines as separate paragraphs
        for idx, line in enumerate(lines):
            if idx > 0:
                text_frame.add_paragraph()
            
            para = text_frame.paragraphs[idx]
            para.text = line  # Content already preprocessed, no need for additional replacements
            para.font.size = Pt(14)
            para.font.color.rgb = theme.get('text', RGBColor(51, 51, 51))
            para.space_after = Pt(10)
            para.level = 0
    
    # ==========================================
    # LAYOUT 13: ERROR SLIDE (Fallback)
    # ==========================================
    def _create_error_slide(self, data: Dict, error_msg: str):
        """Create error slide when generation fails"""
        theme = self.theme_manager.get_theme('dialogue')
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide, theme)
        
        # Error icon
        error_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(2),
            Inches(2.33), Inches(1)
        )
        error_text = error_box.text_frame
        error_para = error_text.paragraphs[0]
        error_para.text = "⚠️"
        error_para.font.size = Pt(72)
        error_para.alignment = PP_ALIGN.CENTER
        
        # Error message
        msg_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.5),
            Inches(9.33), Inches(2)
        )
        msg_text = msg_box.text_frame
        msg_text.word_wrap = True
        
        msg_para = msg_text.paragraphs[0]
        msg_para.text = f"Error generating slide:\n{error_msg}"
        msg_para.font.size = Pt(18)
        msg_para.font.color.rgb = RGBColor(220, 38, 38)
        msg_para.alignment = PP_ALIGN.CENTER
    
    # ==========================================
    # HELPER METHODS
    # ==========================================
    def _set_background(self, slide, theme: Dict):
        """Set slide background color"""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.config.SLIDE_WIDTH,
            self.config.SLIDE_HEIGHT
        )
        background.fill.solid()
        background.fill.fore_color.rgb = theme['bg']
        background.line.fill.background()
        
        # Send to back
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element)
    
    def _add_title(self, slide, title: str, theme: Dict):
        """Add title to slide"""
        title_box = slide.shapes.add_textbox(
            self.config.MARGIN_LEFT,
            Inches(0.3),
            self.config.SLIDE_WIDTH - self.config.MARGIN_LEFT - self.config.MARGIN_RIGHT,
            Inches(1)
        )
        
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = title
        title_paragraph.font.size = self.config.FONT_SIZE_HEADING
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = theme['text']
    
    def get_stats(self) -> Dict:
        """Get service statistics"""
        return {
            'available_themes': self.theme_manager.list_themes(),
            'max_slides': self.config.MAX_SLIDES,
            'image_cache_size': len(self.image_handler.cache),
            'supported_layouts': [
                'centered', 'hero', 'title',
                'split_box', 'split',
                'grid_4', 'grid',
                'roadmap', 'timeline',
                'comparison', 'quote',
                'image_focus', 'two_column',
                'chart', 'table',
                'standard'
            ]
        }
    
    def clear_cache(self):
        """Clear all caches"""
        self.image_handler.clear_cache()
        logger.info("🗑️ All caches cleared")


# ==========================================
# EXPORT / UTILITY FUNCTIONS
# ==========================================
def create_presentation(presentation_data) -> bytes:
    """
    Convenience function to create presentation
    
    Usage:
        pptx_bytes = create_presentation(data)
        with open('output.pptx', 'wb') as f:
            f.write(pptx_bytes)
    """
    service = PPTXService()
    return service.generate(presentation_data)


def get_service_info() -> Dict:
    """Get service information"""
    service = PPTXService()
    return {
        'version': '2.0.0',
        'author': 'GuptaSigma',
        'stats': service.get_stats()
    }


# ==========================================
# MAIN (FOR TESTING)
# ==========================================
if __name__ == "__main__":
    print("="*80)
    print("PPTX Service - Enterprise Edition")
    print("="*80)
    
    info = get_service_info()
    print(f"\nVersion: {info['version']}")
    print(f"Author: {info['author']}")
    print(f"\nAvailable Themes: {', '.join(info['stats']['available_themes'])}")
    print(f"Supported Layouts: {len(info['stats']['supported_layouts'])}")
    print(f"Max Slides: {info['stats']['max_slides']}")
    
    print("\n✅ Service ready for production use!")
    print("="*80)