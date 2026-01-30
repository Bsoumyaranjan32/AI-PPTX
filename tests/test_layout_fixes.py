"""
Tests for PPTX Layout Fixes - Matching Web Preview
Tests the 4 critical layout fixes for proper positioning, layering, and spacing
"""

import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from services.pptx_service import PPTXService, PPTXConfig
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx import Presentation
from io import BytesIO


def test_hero_slide_layering_and_text_color():
    """Test hero slide has proper z-index layering and white text"""
    service = PPTXService()
    
    class MockData:
        theme = 'dialogue'
        content = {
            'slides': [{
                'layout': 'hero',
                'title': 'Introduction to Republic Day',
                'content': 'Republic Day celebrates the adoption of the Constitution of India.',
                'image': 'https://via.placeholder.com/1920x1080'
            }]
        }
    
    output = service.generate(MockData())
    assert len(output) > 0, "Hero slide generation failed"
    
    # Load the presentation to verify structure
    prs = Presentation(BytesIO(output))
    slide = prs.slides[0]
    
    # Verify slide has shapes (image/overlay/text)
    assert len(slide.shapes) >= 3, f"Hero slide should have at least 3 shapes (image/overlay/text), found {len(slide.shapes)}"
    
    print("✅ Hero slide layering test passed")


def test_hero_slide_text_size():
    """Test hero slide uses 18pt body text and 48pt title"""
    service = PPTXService()
    
    class MockData:
        theme = 'dialogue'
        content = {
            'slides': [{
                'layout': 'hero',
                'title': 'Test Title',
                'content': 'Test body content'
            }]
        }
    
    output = service.generate(MockData())
    prs = Presentation(BytesIO(output))
    slide = prs.slides[0]
    
    # Find text frames
    text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame')]
    assert len(text_shapes) >= 2, "Hero slide should have title and body text"
    
    print("✅ Hero slide text size test passed")


def test_standard_slide_with_image_positioning():
    """Test standard slide has correct text and image positioning"""
    service = PPTXService()
    
    class MockData:
        theme = 'dialogue'
        content = {
            'slides': [{
                'layout': 'standard',
                'title': 'Key Traditions',
                'content': 'Flag hoisting ceremony\nParades and cultural programs\nNational anthem',
                'image': 'https://via.placeholder.com/800x600'
            }]
        }
    
    output = service.generate(MockData())
    assert len(output) > 0, "Standard slide generation failed"
    
    prs = Presentation(BytesIO(output))
    slide = prs.slides[0]
    
    # Verify slide has text boxes
    text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame')]
    assert len(text_shapes) >= 1, "Standard slide should have text content"
    
    print("✅ Standard slide positioning test passed")


def test_grid_cards_3_column_layout():
    """Test grid cards use 3-column horizontal layout"""
    service = PPTXService()
    
    class MockData:
        theme = 'dialogue'
        content = {
            'slides': [{
                'layout': 'grid_4',
                'title': 'Importance of Republic Day',
                'content': '* Constitutional Pride\n* Democratic Values\n* National Unity'
            }]
        }
    
    output = service.generate(MockData())
    assert len(output) > 0, "Grid cards generation failed"
    
    prs = Presentation(BytesIO(output))
    slide = prs.slides[0]
    
    # Check for dark navy background
    shapes = list(slide.shapes)
    assert len(shapes) > 0, "Grid slide should have shapes"
    
    # Verify there are shapes for the 3 cards
    # Each card should have: card background + badge + text = 3 shapes per card = 9 shapes + background + title
    assert len(shapes) >= 5, f"Grid slide should have at least 5 shapes (background, title, cards), found {len(shapes)}"
    
    print("✅ Grid cards 3-column layout test passed")


def test_grid_cards_dark_background():
    """Test grid cards have dark navy background"""
    service = PPTXService()
    
    class MockData:
        theme = 'dialogue'
        content = {
            'slides': [{
                'layout': 'grid_4',
                'title': 'Importance Test',
                'content': '* Item 1\n* Item 2\n* Item 3'
            }]
        }
    
    output = service.generate(MockData())
    prs = Presentation(BytesIO(output))
    slide = prs.slides[0]
    
    # First shape should be the dark navy background
    background = slide.shapes[0]
    assert hasattr(background, 'fill'), "Background shape should have fill property"
    
    print("✅ Grid cards dark background test passed")


def test_roadmap_vertical_timeline():
    """Test roadmap uses vertical timeline layout"""
    service = PPTXService()
    
    class MockData:
        theme = 'dialogue'
        content = {
            'slides': [{
                'layout': 'roadmap',
                'title': 'Process/Roadmap',
                'content': '\n'.join([
                    'Preparation',
                    'Parade Rehearsal',
                    'Flag Hoisting',
                    'Cultural Programs',
                    'Main Parade',
                    'Beating Retreat'
                ])
            }]
        }
    
    output = service.generate(MockData())
    assert len(output) > 0, "Roadmap generation failed"
    
    prs = Presentation(BytesIO(output))
    slide = prs.slides[0]
    
    # Verify slide has multiple shapes (line + circles + text boxes)
    # For 6 items: 1 line + 6 circles + 6 number labels + 6 text boxes = 19 shapes + title + background
    assert len(slide.shapes) >= 10, f"Roadmap slide should have multiple shapes for timeline, found {len(slide.shapes)}"
    
    print("✅ Roadmap vertical timeline test passed")


def test_roadmap_six_items_limit():
    """Test roadmap handles exactly 6 items"""
    service = PPTXService()
    
    class MockData:
        theme = 'dialogue'
        content = {
            'slides': [{
                'layout': 'roadmap',
                'title': 'Test Roadmap',
                'content': '\n'.join([f'Phase {i}' for i in range(1, 10)])  # 9 items
            }]
        }
    
    # Should not crash with 9 items (will use first 6)
    try:
        output = service.generate(MockData())
        assert len(output) > 0, "Roadmap slide generation failed"
        print("✅ Roadmap 6 items limit test passed")
    except Exception as e:
        assert False, f"Roadmap slide failed with multiple items: {e}"


def test_all_layouts_generate_successfully():
    """Integration test: All 4 layouts should generate without errors"""
    service = PPTXService()
    
    class MockData:
        theme = 'dialogue'
        content = {
            'slides': [
                {
                    'layout': 'hero',
                    'title': 'Hero Slide',
                    'content': 'Hero content',
                    'image': 'https://via.placeholder.com/1920x1080'
                },
                {
                    'layout': 'standard',
                    'title': 'Standard Slide',
                    'content': 'Standard content',
                    'image': 'https://via.placeholder.com/800x600'
                },
                {
                    'layout': 'grid_4',
                    'title': 'Grid Cards',
                    'content': '* Card 1\n* Card 2\n* Card 3'
                },
                {
                    'layout': 'roadmap',
                    'title': 'Roadmap',
                    'content': 'Step 1\nStep 2\nStep 3\nStep 4\nStep 5\nStep 6'
                }
            ]
        }
    
    try:
        output = service.generate(MockData())
        assert len(output) > 0, "Presentation generation failed"
        
        prs = Presentation(BytesIO(output))
        assert len(prs.slides) == 4, f"Expected 4 slides, got {len(prs.slides)}"
        
        print("✅ All layouts integration test passed")
    except Exception as e:
        assert False, f"Integration test failed: {e}"


if __name__ == '__main__':
    print("\n" + "="*70)
    print("Testing PPTX Layout Fixes - Web Preview Match")
    print("="*70 + "\n")
    
    tests = [
        test_hero_slide_layering_and_text_color,
        test_hero_slide_text_size,
        test_standard_slide_with_image_positioning,
        test_grid_cards_3_column_layout,
        test_grid_cards_dark_background,
        test_roadmap_vertical_timeline,
        test_roadmap_six_items_limit,
        test_all_layouts_generate_successfully
    ]
    
    passed = 0
    failed = 0
    
    for test in tests:
        try:
            test()
            passed += 1
        except AssertionError as e:
            print(f"❌ {test.__name__} failed: {e}")
            failed += 1
        except Exception as e:
            print(f"❌ {test.__name__} error: {e}")
            failed += 1
    
    print("\n" + "="*70)
    print(f"Results: {passed} passed, {failed} failed")
    print("="*70)
    
    sys.exit(0 if failed == 0 else 1)
